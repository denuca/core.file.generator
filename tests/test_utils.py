import io
import logging
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

# Configure logging
logging.basicConfig(level=logging.INFO)

def create_dummy_image(width=100, height=100, color=(255, 0, 0)):
    """Create a dummy image for testing purposes."""
    img = Image.new('RGB', (width, height), color)
    img_bytes = io.BytesIO()
    img.save(img_bytes, format='PNG')
    img_bytes.seek(0)  # Reset the pointer to the beginning
    return img_bytes

def create_pptx_with_images(images):
    """
    Create a PowerPoint presentation with the given images.

    Args:
        images (list): A list of image-like objects (e.g., BytesIO).

    Returns:
        BytesIO: A BytesIO object containing the PPTX file.

    Raises:
        ValueError: If no images are provided.
        Exception: For any other errors during PPTX creation.
    """
    if not images:
        logging.error("No images provided for PPTX creation.")
        raise ValueError("At least one image must be provided.")

    try:
        # Create a new presentation
        pptx_io = BytesIO()
        presentation = Presentation()

        # Add a slide and images to the presentation
        slide_layout = presentation.slide_layouts[5]  # Use a blank slide layout
        slide = presentation.slides.add_slide(slide_layout)

        for image in images:
            # Assuming 'image' is a BytesIO object containing the image data
            left = Inches(1)  # Positioning the image on the slide
            top = Inches(1)
            height = Inches(5)  # Resize the image as needed
            slide.shapes.add_picture(image, left, top, height=height)

        # Save presentation to BytesIO
        presentation.save(pptx_io)
        pptx_io.seek(0)  # Reset pointer to the start for reading later

        logging.info("PPTX created successfully with %d images.", len(images))
        return pptx_io

    except Exception as e:
        logging.exception("Failed to create PPTX: %s", e)
        raise
