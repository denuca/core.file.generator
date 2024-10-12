import logging
import os
import pytest
from core.file_generator.pptx import (
    load_template,
    add_text_box,
    add_image,
    add_new_slide,
    enable_header,
    enable_footer,
    add_slide_number
)
from tests.test_utils import (
    create_dummy_image,
    create_pptx_with_images
)
from io import BytesIO
from pptx import Presentation
from PIL import Image

@pytest.fixture
def ppt_template(tmp_path):
    """Create a temporary PowerPoint file for testing."""
    path = tmp_path / "test.pptx"
    prs = Presentation()
    prs.save(path)
    return path

def test_load_template(ppt_template):
    prs = load_template(ppt_template)
    assert prs is not None

def test_add_text_box(ppt_template):
    prs = load_template(ppt_template)
    slide = add_new_slide(prs, 0)  # Add a title slide
    add_text_box(slide, "Test Text", 1, 1, 5, 1)

    # Adjust the expected shape count based on what is added by `add_new_slide`
    expected_shape_count = 3  # Title shape + Text box

    assert len(slide.shapes) == expected_shape_count  # One for the title slide, one for the text box

def test_enable_header(ppt_template):
    prs = load_template(ppt_template)
    slide = add_new_slide(prs, 1)  # Use a layout that has a header
    enable_header(slide, "Header Text")
    for shape in slide.shapes:
        if shape.has_text_frame:
            assert shape.text_frame.text == "Header Text"
            break

def test_enable_footer(ppt_template):
    prs = load_template(ppt_template)
    slide = add_new_slide(prs, 1)  # Use a layout that has a footer
    enable_footer(slide, "Footer Text")
    for shape in slide.shapes:
        if shape.has_text_frame:
            assert shape.text_frame.text == "Footer Text"
            break

def test_add_slide_number(ppt_template):
    prs = load_template(ppt_template)
    slide = add_new_slide(prs, 1)  # Use a layout that has a footer
    add_slide_number(slide, 1)
    assert any(shape.has_text_frame and shape.text_frame.text == "Slide 1" for shape in slide.shapes)

def test_add_text_box_with_none_size(ppt_template):
    """Test adding a text box with None width and height."""
    prs = load_template(ppt_template)
    slide = add_new_slide(prs, 0)  # Add a title slide
    add_text_box(slide, "Test Text", None, None, None, None)

    # Adjust the expected shape count based on what is added by `add_new_slide`
    expected_shape_count = 3  # Title shape + Text box

    # Check if the text box was added (should be at 0, 0 position with 0 size)
    assert len(slide.shapes) == expected_shape_count  # One for the title slide, one for the text box
    assert slide.shapes[1].shape_type == 14  # Check if it's a text box shape type

def test_create_dummy_image():
    """Test the create_dummy_image function."""
    img_io = create_dummy_image()
    assert img_io is not None
    assert isinstance(img_io, BytesIO)  # Check that the result is a BytesIO object

    # Optionally, you can check the image format or dimensions if needed
    img_io.seek(0)  # Reset the pointer to the start
    img = Image.open(img_io)
    assert img.format == 'PNG'
    assert img.size == (100, 100)  # Confirm dimensions of the dummy image

def test_create_pptx_with_images():
    """Test the PowerPoint creation function with images."""
    dummy_image = create_dummy_image()  # Generate a dummy image
    pptx_io = create_pptx_with_images([dummy_image])  # Use the utility function

    # Verify that a PPTX file is created
    assert pptx_io is not None
    assert isinstance(pptx_io, BytesIO)

    # Optionally, you can validate the contents of the PPTX file
    pptx_io.seek(0)  # Reset the pointer to the start
    presentation = Presentation(pptx_io)

    # Verify the number of slides
    assert len(presentation.slides) == 1  # Check that there is one slide
    slide = presentation.slides[0]

    # Verify that only the expected shapes are present
    image_shapes = [shape for shape in slide.shapes if shape.shape_type == 13]  # Only look for image shapes
    assert len(image_shapes) == 1  # Check that there is exactly one image shape

    # Log the shape types for debugging
    for shape in slide.shapes:
        logging.info(f"Shape type: {shape.shape_type}")

    # If you need to ensure the image is in the right position, you can add additional checks

def test_add_image_with_none_size(ppt_template):
    """Test adding an image with None width and height."""
    prs = load_template(ppt_template)
    slide = add_new_slide(prs, 0)  # Add a title slide
    # Create a dummy image file
    image_io = create_dummy_image()

    # Test adding an image with None for width and height
    add_image(slide, image_io, 1, 1, None, None)

    # Adjust the expected shape count based on what is added by `add_new_slide`
    expected_shape_count = 3  # Title shape + Text box + Image

    # Check if the image was added
    assert len(slide.shapes) == expected_shape_count  # One for the title slide, one for the text box, one for the image
    image_shapes = [shape for shape in slide.shapes if shape.shape_type == 13]  # Only look for image shapes
    assert len(image_shapes) == 1  # Check that there is exactly one image shape

    # Log the shape types for debugging
    for shape in slide.shapes:
        logging.info(f"Shape type: {shape.shape_type}")

def test_add_image_with_dimensions(ppt_template):
    """Test adding an image to a slide with specified dimensions."""
    prs = load_template(ppt_template)
    slide = add_new_slide(prs, 0)  # Add a title slide
    # Create a dummy image file
    image_io = create_dummy_image()

    # Test adding an image with specified width and height
    add_image(slide, image_io, 1, 1, 2, 2)  # Add image with width and height

    # Adjust the expected shape count based on what is added by `add_new_slide`
    expected_shape_count = 3  # Title shape + Text box + Image

    # Check if the image was added
    assert len(slide.shapes) == expected_shape_count  # One for the title slide, one for the text box, one for the image

def test_add_text_box_with_dimensions(ppt_template):
    """Test adding a text box to a slide with specified dimensions."""
    prs = load_template(ppt_template)
    slide = add_new_slide(prs, 0)  # Add a title slide
    add_text_box(slide, "Test Text", 1, 1, 5, 1)  # Add text box with specified dimensions

    # Adjust the expected shape count based on what is added by `add_new_slide`
    expected_shape_count = 3  # Title shape + Text box

    # Check if the text box was added
    assert len(slide.shapes) == expected_shape_count  # One for the title slide, one for the text box
    assert slide.shapes[1].shape_type == 14  # Check if it's a text box shape type
