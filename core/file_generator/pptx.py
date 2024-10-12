import logging
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def load_template(template_path):
    """Load a PowerPoint template."""
    try:
        return Presentation(template_path)
    except Exception as e:
        logger.error(f"Error loading template {template_path}: {e}")
        raise

def get_inches(left=None, top=None, width=None, height=None):
    """Convert provided dimensions to inches, handling None values."""
    try:
        left_inch = Inches(left) if left is not None else Inches(0)
        top_inch = Inches(top) if top is not None else Inches(0)
        width_inch = Inches(width) if width is not None else Inches(0)
        height_inch = Inches(height) if height is not None else Inches(0)

        logger.debug(f"Converted to inches: left={left_inch}, top={top_inch}, width={width_inch}, height={height_inch}")
        return left_inch, top_inch, width_inch, height_inch
    except Exception as e:
        logger.error(f"Error converting to inches: {e}")
        raise
from pptx import Presentation

def create_ppt():
    """Create a new PowerPoint presentation object."""
    return Presentation()

def add_title_slide(ppt, title, content):
    """Add a title slide with title and content to the presentation."""
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])  # Title Slide
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    content_placeholder.text = content
    return slide

def add_text_to_slide(slide, placeholder_name, text):
    """Add text to a slide using a placeholder name."""
    try:
        for shape in slide.shapes:
            if shape.name == placeholder_name:
                shape.text = text
                logger.info(f"Added text to placeholder '{placeholder_name}': {text}")
                break
    except Exception as e:
        logger.error(f"Error adding text to slide: {e}")
        raise

def add_text_box(slide, text, left, top, width, height):
    """Add a text box to the slide at specified position."""
    try:
        left_inch, top_inch, width_inch, height_inch = get_inches(left, top, width, height)

        textbox = slide.shapes.add_textbox(
            left_inch,
            top_inch,
            width_inch,
            height_inch
        )
        text_frame = textbox.text_frame
        text_frame.text = text
        logger.info(f"Added text box: '{text}' at position ({left}, {top}) with size ({width}, {height})")
    except Exception as e:
        logger.error(f"Error adding text box: {e}")
        raise

def add_image(slide, image_path, left, top, width=None, height=None):
    """Add an image to the slide at specified position."""
    try:
        left_inch, top_inch, width_inch, height_inch = get_inches(left, top, width if width else 0, height if height else 0)
        
        logger.info(f"Adding image '{image_path}' at position ({left}, {top}) with size ({width}, {height})")

        slide.shapes.add_picture(
            image_path,
            left_inch,
            top_inch,
            width=width_inch if width else None,
            height=height_inch if height else None
        )
    except Exception as e:
        logger.error(f"Error adding image to slide: {e}")
        raise

def add_new_slide(prs, layout):
    """Add a new slide with the specified layout."""
    try:
        slide_layout = prs.slide_layouts[layout]  # layout is an index (e.g., 0 for title slide)
        slide = prs.slides.add_slide(slide_layout)
        logger.info(f"Added new slide with layout index {layout}")
        return slide
    except Exception as e:
        logger.error(f"Error adding new slide: {e}")
        raise

def enable_header(slide, header_text):
    """Enable and set text for the header of the slide."""
    try:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text == "":
                shape.text_frame.text = header_text
                logger.info(f"Enabled header with text: {header_text}")
                break
    except Exception as e:
        logger.error(f"Error enabling header: {e}")
        raise

def enable_footer(slide, footer_text):
    """Enable and set text for the footer of the slide."""
    try:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text == "":
                shape.text_frame.text = footer_text
                logger.info(f"Enabled footer with text: {footer_text}")
                break
    except Exception as e:
        logger.error(f"Error enabling footer: {e}")
        raise

def add_slide_number(slide, slide_number):
    """Add pagination (slide number) to the footer of the slide."""
    try:
        left_inch, top_inch, width_inch, height_inch = get_inches(8.5, 5.5, 1, 0.5)

        footer_shape = slide.shapes.add_textbox(
            left_inch,
            top_inch,
            width_inch,
            height_inch
        )
        footer_shape.text_frame.text = f"Slide {slide_number}"
        logger.info(f"Added slide number: {slide_number} at position (8.5, 5.5)")
    except Exception as e:
        logger.error(f"Error adding slide number: {e}")
        raise

def save_presentation(presentation, output_path):
    """Save the PowerPoint presentation."""
    try:
        presentation.save(output_path)
        logger.info(f"Saved presentation to {output_path}")
    except Exception as e:
        logger.error(f"Error saving presentation to {output_path}: {e}")
        raise
